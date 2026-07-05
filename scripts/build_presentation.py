#!/usr/bin/env python3
"""
Build the AI Final Examination presentation from the course template.

Usage:
    python scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "AI Final Examination Presentation Template.pptx"
OUTPUT = ROOT / "ppt" / "Student_Performance_Prediction_Agent.pptx"
OUTPUTS = ROOT / "outputs"

# --- Customize before presenting ---
STUDENT_NAME = "[Your Name]"
SEMESTER = "[Spring 2026]"
GROUP = "[Group Number]"


def set_paragraphs(text_frame, lines: list[str], font_size: int = 14) -> None:
    """Replace all paragraphs in a text frame with new lines."""
    tf = text_frame
    # Clear existing paragraphs except first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._element
        p.getparent().remove(p)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"


def set_table_cell(table, row: int, col: int, text: str, bold: bool = False) -> None:
    cell = table.cell(row, col)
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(11)
            run.font.name = "Calibri"
            run.font.bold = bold


def replace_pictures(slide, image_paths: list[Path]) -> None:
    """Remove existing pictures on slide and insert project images at similar positions."""
    pic_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    positions = [(s.left, s.top, s.width, s.height) for s in pic_shapes]

    for shape in pic_shapes:
        el = shape._element
        el.getparent().remove(el)

    # Default layout if template pictures were removed
    if not positions:
        positions = [
            (Inches(0.5), Inches(2.0), Inches(4.5), Inches(3.0)),
            (Inches(5.2), Inches(2.0), Inches(4.3), Inches(3.0)),
        ]

    for path, (left, top, width, height) in zip(image_paths, positions):
        if path.exists():
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_picture_if_room(slide, path: Path, left, top, width, height) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def create_architecture_diagram(path: Path) -> None:
    """Create a complete agent architecture diagram with all components and feedback loops."""
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

    fig, ax = plt.subplots(figsize=(14, 7))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 7)
    ax.axis("off")
    ax.set_facecolor("#fafafa")

    def box(x, y, w, h, label, color, fontsize=8.5):
        patch = FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.04,rounding_size=0.12",
            linewidth=1.8, edgecolor="#2c3e50", facecolor=color, alpha=0.92,
        )
        ax.add_patch(patch)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center",
                fontsize=fontsize, color="white", fontweight="bold", linespacing=1.25)

    def arrow(x1, y1, x2, y2, style="-|>", color="#2c3e50", lw=1.8, connection="arc3"):
        ax.add_patch(FancyArrowPatch(
            (x1, y1), (x2, y2),
            arrowstyle=style, color=color, lw=lw,
            connectionstyle=connection, mutation_scale=14,
        ))

    # Title
    ax.text(7, 6.55, "Student Performance Prediction Agent — Full Architecture",
            ha="center", fontsize=15, fontweight="bold", color="#2c3e50")
    ax.text(7, 6.15, "Perceive → Learn → Reason → Act  (Russell & Norvig Intelligent Agent)",
            ha="center", fontsize=10, color="#7f8c8d", style="italic")

    # --- Main pipeline (top row) ---
    bw, bh = 1.55, 1.05
    y_main = 4.35
    pipeline = [
        (0.4, "Sensors\n(SIS · LMS ·\nAttendance)"),
        (2.25, "Percepts\n(Student\nFeature Vector)"),
        (4.1, "Preprocessor\n(Scale · Encode)"),
        (5.95, "Random Forest\n(Learning\nComponent)"),
        (7.8, "Decision\nEngine\n(Risk + Rules)"),
    ]
    colors_main = ["#2980b9", "#8e44ad", "#16a085", "#d35400", "#c0392b"]
    for (x, label), color in zip(pipeline, colors_main):
        box(x, y_main, bw, bh, label, color)

    # Actuators (right column)
    actuators = [
        (9.85, 5.15, "Teacher\nDashboard"),
        (9.85, 4.0, "Notifications\n(Counselor · Parent)"),
        (9.85, 2.85, "Recommendations\n(Tutoring · Study Plan)"),
    ]
    colors_act = ["#27ae60", "#2ecc71", "#1e8449"]
    for (x, y, label), color in zip(actuators, colors_act):
        box(x, y, bw, bh, label, color)

    # Support components (bottom row)
    box(3.2, 1.35, 2.2, 1.05, "Knowledge Base\n(Thresholds ·\nIntervention Rules)", "#34495e")
    box(6.1, 1.35, 2.2, 1.05, "Memory\n(Prediction History\nper Student)", "#5d6d7e")

    # Environment label
    box(0.35, 5.75, 1.55, 0.55, "Environment\n(School)", "#95a5a6", fontsize=8)
    arrow(1.1, 5.75, 1.1, 5.4)

    # Main flow arrows
    arrow(1.95, y_main + bh / 2, 2.25, y_main + bh / 2)
    arrow(3.8, y_main + bh / 2, 4.1, y_main + bh / 2)
    arrow(5.65, y_main + bh / 2, 5.95, y_main + bh / 2)
    arrow(7.5, y_main + bh / 2, 7.8, y_main + bh / 2)
    arrow(9.35, y_main + bh / 2, 9.85, y_main + bh / 2 + 0.35)

    # Decision engine → actuators
    arrow(8.575, y_main, 9.85, 5.7, connection="arc3,rad=0.15")
    arrow(8.575, y_main + 0.2, 9.85, 4.55, connection="arc3,rad=0.1")
    arrow(8.575, y_main, 9.85, 3.4, connection="arc3,rad=-0.15")

    # Knowledge base → decision engine
    arrow(4.3, 2.4, 8.2, 4.35, connection="arc3,rad=-0.2", color="#34495e")
    ax.text(5.8, 3.35, "rules", fontsize=8, color="#34495e", style="italic")

    # Random forest ↔ knowledge (training context)
    arrow(6.72, 4.35, 5.5, 2.4, connection="arc3,rad=0.25", color="#7f8c8d", style="-|>")
    ax.text(6.0, 3.55, "beliefs", fontsize=8, color="#7f8c8d", style="italic")

    # Memory ↔ decision engine
    arrow(7.2, 2.4, 8.0, 4.35, connection="arc3,rad=0.2", color="#5d6d7e")
    arrow(8.0, 4.35, 7.2, 2.4, connection="arc3,rad=0.35", color="#5d6d7e", style="-|>")

    # Feedback loop: dashboard → memory
    arrow(10.6, 2.85, 7.2, 2.4, connection="arc3,rad=-0.35", color="#e67e22", lw=2)
    ax.text(9.2, 2.15, "feedback loop", fontsize=8, color="#e67e22", fontweight="bold")

    # Legend
    legend_y = 0.35
    ax.text(0.4, legend_y, "● Sensors/Actuators   ● Cognition   ● Learning   ● Support",
            fontsize=8.5, color="#566573")

    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", pad_inches=0.25)
    plt.close(fig)


def build() -> Path:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))

    # --- Slide 1: Agent ---
    s1 = prs.slides[0]
    shapes = [sh for sh in s1.shapes if sh.has_text_frame]
    if len(shapes) >= 2:
        shapes[0].text_frame.text = "Student Performance Prediction Agent"
        for run in shapes[0].text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.bold = True

        set_paragraphs(
            shapes[1].text_frame,
            [
                "Agent name:",
                '"Student Performance Prediction Agent (SPPA)"',
                "",
                "Goal:",
                "Predict student performance (Low / Medium / High) and recommend "
                "early interventions before final examinations.",
                "",
                f"{STUDENT_NAME}  |  {SEMESTER}  |  Group {GROUP}",
                "Artificial Intelligence — Final Examination",
            ],
            font_size=13,
        )

    # --- Slide 2: Environment ---
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if "environment/context" in txt or "Russell" in txt:
                if "environment/context" in txt:
                    shape.text_frame.text = "Russell & Norvig Environment Analysis"
                    for run in shape.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                elif "Russell" in txt:
                    set_paragraphs(
                        shape.text_frame,
                        [
                            "Educational ecosystem: schools, classrooms, homes, examination periods",
                            "Agent perceives partial student records and acts via teacher dashboard",
                        ],
                        font_size=12,
                    )
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = shape.table
            rows = [
                ("Observable", "Partially Observable ✓", "Fully Observable",
                 "Hidden motivation, home stress, undiagnosed learning needs"),
                ("Deterministic vs Stochastic", "Stochastic ✓", "Deterministic",
                 "Random Forest outputs probabilities; outcomes are uncertain"),
                ("Episodic vs Sequential", "Sequential ✓", "Episodic",
                 "G1/G2 history affects G3; interventions shape future state"),
                ("Static vs Dynamic", "Dynamic ✓", "Static",
                 "Grades and absences change throughout the semester"),
                ("Discrete vs Continuous", "Discrete ✓", "Continuous",
                 "Decisions: Low / Medium / High performance categories"),
                ("Single vs Multi-agent", "Multi-agent ✓", "Single",
                 "Teachers, students, parents, counselors interact"),
            ]
            for i, (prop, chosen, alt, reason) in enumerate(rows, start=1):
                set_table_cell(tbl, i, 0, prop, bold=True)
                set_table_cell(tbl, i, 1, chosen)
                set_table_cell(tbl, i, 2, alt)
                set_table_cell(tbl, i, 3, reason)

    # --- Slide 3: Dataset ---
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame and "One slide" in shape.text_frame.text:
            shape.text_frame.text = "Dataset & Sensors"
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(24)
                run.font.bold = True
        if shape.has_text_frame and "Where did it come from" in shape.text_frame.text:
            set_paragraphs(
                shape.text_frame,
                [
                    "Source: UCI ML Repository — Student Performance (Cortez & Silva, 2008)",
                    "Description: 1,044 records (395 Math + 649 Portuguese), 33 features",
                    "Sensors: SIS grades, attendance logs, enrollment forms, study surveys",
                    "Target: G3 final grade (0–20) → Low / Medium / High",
                    "Issues: 22% Low class imbalance, absences outliers (max 93), self-reported fields",
                ],
                font_size=12,
            )
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = shape.table
            set_table_cell(tbl, 0, 0, "Sample fields", bold=True)
            set_table_cell(tbl, 0, 1, "age, sex, studytime", bold=True)
            set_table_cell(tbl, 0, 2, "failures, absences, G1, G2", bold=True)
            set_table_cell(tbl, 0, 3, "G3 / Category", bold=True)
            set_table_cell(tbl, 1, 0, "Sample rows")
            set_table_cell(tbl, 1, 1, "15,F,2 | 18,M,1 | 17,F,4")
            set_table_cell(tbl, 1, 2, "failures, absences, G1, G2")
            set_table_cell(tbl, 1, 3, "10/Med · 6/Low · 15/High")

    # --- Slide 4: EDA ---
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.has_text_frame and "inspection" in shape.text_frame.text.lower():
            shape.text_frame.text = "Exploratory Data Analysis"
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(24)
                run.font.bold = True
        if shape.has_text_frame and "covariance" in shape.text_frame.text.lower():
            set_paragraphs(
                shape.text_frame,
                [
                    "G2 ↔ G3 correlation: r ≈ 0.90 (strongest predictor)",
                    "G1 ↔ G3: r ≈ 0.80  |  failures ↔ G3: r ≈ -0.35",
                    "Mean G3 = 11.9 (σ = 3.2)  |  0 missing values",
                    "Categories: Low 22% · Medium 50% · High 28%",
                ],
                font_size=12,
            )
    replace_pictures(
        s4,
        [
            OUTPUTS / "correlation_heatmap.png",
            OUTPUTS / "class_distribution.png",
        ],
    )

    # --- Slide 5: Model ---
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.has_text_frame and "model/algorithm" in shape.text_frame.text.lower():
            shape.text_frame.text = "Machine Learning — Learning Component"
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(24)
                run.font.bold = True
        if shape.has_text_frame and "supervised" in shape.text_frame.text.lower():
            set_paragraphs(
                shape.text_frame,
                [
                    "Type: Supervised Multi-Class Classification",
                    "Pipeline: Features → StandardScaler + OneHotEncoder → Random Forest",
                    "Best model: Random Forest (200 trees, depth 12, balanced weights)",
                    "Compared 7 models — RF wins: F1 = 86.15%, Accuracy = 86.12%",
                    "Runners-up: Logistic Regression (85.16%), XGBoost (84.68%)",
                ],
                font_size=12,
            )
    add_picture_if_room(
        s5,
        OUTPUTS / "model_comparison.png",
        Inches(5.0),
        Inches(1.8),
        Inches(4.5),
        Inches(3.2),
    )

    # --- Slide 6: Results ---
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame and "Results" in shape.text_frame.text:
            shape.text_frame.text = "Evaluation Results — Random Forest"
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(24)
                run.font.bold = True
        if shape.has_text_frame and "Accuracy" in shape.text_frame.text:
            set_paragraphs(
                shape.text_frame,
                [
                    "Accuracy: 86.12%  |  Precision: 86.36%  |  Recall: 86.12%",
                    "F1 Score: 86.15%  |  ROC AUC: 94.61%",
                    "CV Accuracy: 84.67% ± 1.95%  |  Train/Test: 835 / 209",
                    "Top predictor: G2 (2nd period grade)",
                    "Training: 0.43s  |  Prediction: 34ms per batch",
                ],
                font_size=12,
            )
    add_picture_if_room(
        s6,
        OUTPUTS / "confusion_matrix.png",
        Inches(0.4),
        Inches(2.0),
        Inches(3.8),
        Inches(2.8),
    )
    add_picture_if_room(
        s6,
        OUTPUTS / "feature_importance.png",
        Inches(4.5),
        Inches(2.0),
        Inches(4.8),
        Inches(2.8),
    )
    add_picture_if_room(
        s6,
        OUTPUTS / "roc_curve.png",
        Inches(9.5),
        Inches(2.0),
        Inches(3.5),
        Inches(2.8),
    )

    # --- Slide 7: Agent Integration ---
    arch_path = OUTPUTS / "agent_architecture.png"
    create_architecture_diagram(arch_path)

    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame and "Final slide" in shape.text_frame.text:
            shape.text_frame.text = "Agent Integration — From Prediction to Action"
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(22)
                run.font.bold = True
        if shape.has_text_frame and "How could" in shape.text_frame.text:
            set_paragraphs(
                shape.text_frame,
                [
                    "ML model = Learning Component inside the intelligent agent",
                    "Flow: Percepts → Preprocess → RF probabilities → Risk engine → Actions",
                    "Dashboard scores 1,044 students: 233 HIGH · 525 MEDIUM · 286 LOW priority",
                    "Actions: counseling, tutoring, parent alerts, teacher dashboard CSV",
                    "Components: Sensors, Knowledge Base, Decision Engine, Memory, Actuators",
                    "Limitations: Portuguese data only, ~14% error, privacy & bias concerns",
                    "Future: SHAP explainability, real-time analytics, LLM parent reports",
                    "",
                    "Thank you — Questions welcome.",
                ],
                font_size=12,
            )
    add_picture_if_room(
        s7,
        arch_path,
        Inches(0.25),
        Inches(3.35),
        Inches(9.3),
        Inches(2.1),
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Presentation saved to: {out}")
